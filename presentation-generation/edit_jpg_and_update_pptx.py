#!/usr/bin/env python3
"""Apply simple JPG edits and place the result into a PowerPoint deck.

This helper is intended to sit after the existing generation scripts:

    export_slide_images.ps1 -> build_presentation_deck.ps1 -> this script

Image edits are driven by a small JSON file so the changes can be repeated.
PowerPoint insertion uses python-pptx and supports replacing an existing slide,
inserting after a slide, or appending a new full-image slide.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


DEFAULT_FONT_PATHS = [
    Path(r"C:\Windows\Fonts\arial.ttf"),
    Path(r"C:\Windows\Fonts\calibri.ttf"),
]
DEFAULT_BOLD_FONT_PATHS = [
    Path(r"C:\Windows\Fonts\arialbd.ttf"),
    Path(r"C:\Windows\Fonts\calibrib.ttf"),
]


def parse_color(value: str | list[int] | tuple[int, ...] | None, default: str) -> tuple[int, int, int, int]:
    if value is None:
        value = default
    if isinstance(value, (list, tuple)):
        if len(value) == 3:
            return int(value[0]), int(value[1]), int(value[2]), 255
        if len(value) == 4:
            return int(value[0]), int(value[1]), int(value[2]), int(value[3])
        raise ValueError(f"Color arrays must have 3 or 4 values: {value!r}")
    value = value.strip()
    named = {
        "white": "#ffffff",
        "black": "#000000",
        "navy": "#082659",
        "blue": "#00539b",
        "red": "#c00000",
    }
    value = named.get(value.lower(), value)
    if value.startswith("#"):
        value = value[1:]
    if len(value) == 6:
        return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16), 255
    if len(value) == 8:
        return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16), int(value[6:8], 16)
    raise ValueError(f"Unsupported color value: {value!r}")


def load_font(size: int, bold: bool = False, font_path: str | None = None) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [Path(font_path)] if font_path else (DEFAULT_BOLD_FONT_PATHS if bold else DEFAULT_FONT_PATHS)
    for path in candidates:
        if path and path.exists():
            return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


def box_from_op(operation: dict[str, Any]) -> tuple[int, int, int, int]:
    raw = operation.get("box")
    if not isinstance(raw, list) or len(raw) != 4:
        raise ValueError(f"Operation requires box [x, y, width, height]: {operation!r}")
    x, y, w, h = [int(v) for v in raw]
    return x, y, x + w, y + h


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.ImageFont, max_width: int) -> list[str]:
    words = text.split()
    lines: list[str] = []
    current = ""
    for word in words:
        test = word if not current else f"{current} {word}"
        if draw.textlength(test, font=font) <= max_width:
            current = test
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def draw_text_operation(draw: ImageDraw.ImageDraw, operation: dict[str, Any]) -> None:
    text = str(operation["text"])
    x, y = [int(v) for v in operation.get("position", [0, 0])]
    size = int(operation.get("font_size", 24))
    font = load_font(size=size, bold=bool(operation.get("bold", False)), font_path=operation.get("font"))
    fill = parse_color(operation.get("fill"), "#000000")
    anchor = operation.get("anchor")
    max_width = operation.get("max_width")
    line_spacing = int(operation.get("line_spacing", 4))

    if max_width:
        for line in wrap_text(draw, text, font, int(max_width)):
            draw.text((x, y), line, font=font, fill=fill, anchor=anchor)
            y += size + line_spacing
    else:
        draw.text((x, y), text, font=font, fill=fill, anchor=anchor)


def paste_image(base: Image.Image, operation: dict[str, Any], spec_dir: Path) -> None:
    image_path = Path(operation["image"])
    if not image_path.is_absolute():
        image_path = spec_dir / image_path
    overlay = Image.open(image_path).convert("RGBA")
    x1, y1, x2, y2 = box_from_op(operation)
    target_w = x2 - x1
    target_h = y2 - y1
    fit = operation.get("fit", "contain")

    if fit == "stretch":
        resized = overlay.resize((target_w, target_h), Image.Resampling.LANCZOS)
        base.paste(resized, (x1, y1), resized)
        return

    src_w, src_h = overlay.size
    scale = min(target_w / src_w, target_h / src_h) if fit == "contain" else max(target_w / src_w, target_h / src_h)
    new_size = (round(src_w * scale), round(src_h * scale))
    resized = overlay.resize(new_size, Image.Resampling.LANCZOS)
    offset_x = x1 + (target_w - new_size[0]) // 2
    offset_y = y1 + (target_h - new_size[1]) // 2
    base.paste(resized, (offset_x, offset_y), resized)


def apply_image_operations(source_image: Path, output_image: Path, spec_path: Path) -> Path:
    with spec_path.open("r", encoding="utf-8") as handle:
        spec = json.load(handle)

    image = Image.open(source_image).convert("RGBA")
    draw = ImageDraw.Draw(image)
    spec_dir = spec_path.parent

    for operation in spec.get("operations", []):
        op_type = operation.get("type")
        if op_type == "crop":
            raw_box = operation.get("box")
            if not isinstance(raw_box, list) or len(raw_box) != 4:
                raise ValueError("crop requires box [left, top, right, bottom]")
            image = image.crop(tuple(int(v) for v in raw_box))
            draw = ImageDraw.Draw(image)
        elif op_type == "resize":
            raw_size = operation.get("size")
            if not isinstance(raw_size, list) or len(raw_size) != 2:
                raise ValueError("resize requires size [width, height]")
            image = image.resize(tuple(int(v) for v in raw_size), Image.Resampling.LANCZOS)
            draw = ImageDraw.Draw(image)
        elif op_type in {"cover", "rectangle"}:
            fill = parse_color(operation.get("fill"), "#ffffff") if op_type == "cover" else None
            outline = parse_color(operation.get("outline"), "#000000") if operation.get("outline") else None
            width = int(operation.get("width", 1))
            radius = int(operation.get("radius", 0))
            box = box_from_op(operation)
            if radius:
                draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)
            else:
                draw.rectangle(box, fill=fill, outline=outline, width=width)
        elif op_type == "text":
            draw_text_operation(draw, operation)
        elif op_type == "paste":
            paste_image(image, operation, spec_dir)
        else:
            raise ValueError(f"Unsupported operation type: {op_type!r}")

    output_image.parent.mkdir(parents=True, exist_ok=True)
    save_kwargs: dict[str, Any] = {}
    if output_image.suffix.lower() in {".jpg", ".jpeg"}:
        image = image.convert("RGB")
        save_kwargs["quality"] = int(spec.get("quality", 95))
    image.save(output_image, **save_kwargs)
    return output_image


def blank_slide_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if len(layout.placeholders) == 0:
            return layout
    return prs.slide_layouts[-1]


def clear_slide(slide) -> None:
    shape_elements = list(slide.shapes._spTree)
    for element in shape_elements:
        if element.tag.endswith("}nvGrpSpPr") or element.tag.endswith("}grpSpPr"):
            continue
        slide.shapes._spTree.remove(element)


def move_slide(prs: Presentation, old_index: int, new_index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[old_index]
    slide_id_list.remove(slide_id)
    slide_id_list.insert(new_index, slide_id)


def add_full_slide_picture(slide, image_path: Path, prs: Presentation, fit: str) -> None:
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    if fit == "stretch":
        slide.shapes.add_picture(str(image_path), 0, 0, width=slide_w, height=slide_h)
        return

    with Image.open(image_path) as image:
        img_w, img_h = image.size

    scale = min(slide_w / img_w, slide_h / img_h) if fit == "contain" else max(slide_w / img_w, slide_h / img_h)
    width = round(img_w * scale)
    height = round(img_h * scale)
    left = round((slide_w - width) / 2)
    top = round((slide_h - height) / 2)
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def update_presentation(
    presentation_path: Path,
    image_path: Path,
    output_path: Path,
    slide_number: int | None,
    mode: str,
    fit: str,
) -> Path:
    prs = Presentation(str(presentation_path))

    if mode == "replace":
        if slide_number is None:
            raise ValueError("--slide-number is required for replace mode")
        if slide_number < 1 or slide_number > len(prs.slides):
            raise ValueError(f"Slide number {slide_number} is outside the deck range 1-{len(prs.slides)}")
        slide = prs.slides[slide_number - 1]
        clear_slide(slide)
    else:
        slide = prs.slides.add_slide(blank_slide_layout(prs))
        if mode == "insert-after":
            if slide_number is None:
                raise ValueError("--slide-number is required for insert-after mode")
            if slide_number < 0 or slide_number > len(prs.slides) - 1:
                raise ValueError(f"Insert-after slide number must be between 0 and {len(prs.slides) - 1}")
            move_slide(prs, len(prs.slides) - 1, slide_number)

    add_full_slide_picture(slide, image_path, prs, fit)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Edit a JPG/PNG and insert it into a PowerPoint presentation.")
    parser.add_argument("--source-image", type=Path, help="Input JPG/PNG to edit. Required when --edit-spec is used.")
    parser.add_argument("--edit-spec", type=Path, help="JSON file with image edit operations.")
    parser.add_argument("--output-image", type=Path, help="Where to write the edited image.")
    parser.add_argument("--presentation", type=Path, help="PowerPoint file to update.")
    parser.add_argument("--output-presentation", type=Path, help="Where to write the updated PowerPoint file.")
    parser.add_argument("--slide-number", type=int, help="1-based slide number for replace/insert-after.")
    parser.add_argument("--insert-mode", choices=["replace", "insert-after", "append"], default="replace")
    parser.add_argument("--picture-fit", choices=["contain", "cover", "stretch"], default="contain")
    return parser


def main() -> int:
    args = build_parser().parse_args()

    image_to_insert = args.output_image or args.source_image
    if args.edit_spec:
        if not args.source_image or not args.output_image:
            raise SystemExit("--source-image and --output-image are required with --edit-spec")
        image_to_insert = apply_image_operations(args.source_image, args.output_image, args.edit_spec)
        print(f"Created edited image: {image_to_insert}")

    if args.presentation:
        if not image_to_insert:
            raise SystemExit("Provide --output-image or --source-image before updating a presentation")
        if not args.output_presentation:
            raise SystemExit("--output-presentation is required with --presentation")
        updated = update_presentation(
            presentation_path=args.presentation,
            image_path=image_to_insert,
            output_path=args.output_presentation,
            slide_number=args.slide_number,
            mode=args.insert_mode,
            fit=args.picture_fit,
        )
        print(f"Created updated presentation: {updated}")

    if not args.edit_spec and not args.presentation:
        raise SystemExit("Nothing to do. Provide --edit-spec and/or --presentation.")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
