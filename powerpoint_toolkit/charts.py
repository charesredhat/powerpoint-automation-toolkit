"""Chart creation helpers."""

from __future__ import annotations

from typing import Dict, List, Optional, Sequence, Union

from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Public mapping of friendly names to XL_CHART_TYPE members
CHART_TYPES: Dict[str, XL_CHART_TYPE] = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}


class ChartBuilder:
    """Fluent builder for adding charts to a slide.

    Parameters
    ----------
    chart_type:
        A friendly chart-type name (see :data:`CHART_TYPES`) or an
        :class:`~pptx.enum.chart.XL_CHART_TYPE` value directly.

    Examples
    --------
    Add a clustered column chart::

        from powerpoint_toolkit import Presentation
        from powerpoint_toolkit.charts import ChartBuilder

        prs = Presentation()
        slide = prs.add_slide(layout="blank")

        chart = (
            ChartBuilder("column")
            .categories(["Q1", "Q2", "Q3", "Q4"])
            .add_series("Revenue", [10, 20, 15, 30])
            .add_series("Profit",  [3,  8,  5,  12])
            .title("Quarterly Performance")
            .position(left=1, top=1.5, width=8, height=4.5)
            .build(slide)
        )
    """

    def __init__(self, chart_type: Union[str, XL_CHART_TYPE] = "column") -> None:
        if isinstance(chart_type, str):
            key = chart_type.lower().replace(" ", "_").replace("-", "_")
            if key not in CHART_TYPES:
                available = ", ".join(sorted(CHART_TYPES))
                raise ValueError(
                    f"Unknown chart type {chart_type!r}.  Available: {available}"
                )
            self._chart_type: XL_CHART_TYPE = CHART_TYPES[key]
        else:
            self._chart_type = chart_type

        self._categories: List[str] = []
        self._series: List[tuple] = []  # list of (name, values)
        self._title: Optional[str] = None
        self._has_legend: bool = True
        self._has_data_labels: bool = False

        # Position defaults (inches)
        self._left: float = 1.0
        self._top: float = 1.5
        self._width: float = 8.0
        self._height: float = 4.5

    # ------------------------------------------------------------------
    # Fluent configuration
    # ------------------------------------------------------------------

    def categories(self, labels: Sequence[str]) -> "ChartBuilder":
        """Set the category labels (x-axis for most chart types)."""
        self._categories = list(labels)
        return self

    def add_series(
        self, name: str, values: Sequence[Union[int, float]]
    ) -> "ChartBuilder":
        """Add a data series.

        Parameters
        ----------
        name:
            Series name (shown in the legend).
        values:
            Numeric values, one per category.
        """
        self._series.append((name, list(values)))
        return self

    def title(self, text: str) -> "ChartBuilder":
        """Set the chart title."""
        self._title = text
        return self

    def legend(self, show: bool = True) -> "ChartBuilder":
        """Show or hide the chart legend."""
        self._has_legend = show
        return self

    def data_labels(self, show: bool = True) -> "ChartBuilder":
        """Show or hide data labels on chart series."""
        self._has_data_labels = show
        return self

    def position(
        self,
        left: float = 1.0,
        top: float = 1.5,
        width: float = 8.0,
        height: float = 4.5,
    ) -> "ChartBuilder":
        """Set the chart position and size (all values in inches)."""
        self._left = left
        self._top = top
        self._width = width
        self._height = height
        return self

    # ------------------------------------------------------------------
    # Build
    # ------------------------------------------------------------------

    def build(self, slide: "Slide") -> object:  # noqa: F821
        """Add the chart to *slide* and return the underlying chart object.

        Parameters
        ----------
        slide:
            A :class:`~powerpoint_toolkit.slides.Slide` instance (or any
            object exposing a ``.pptx_slide`` attribute).

        Returns
        -------
        pptx.chart.chart.Chart
            The underlying python-pptx chart object, which can be further
            customised if needed.
        """
        chart_data = CategoryChartData()

        if self._categories:
            chart_data.categories = self._categories

        for name, values in self._series:
            chart_data.add_series(name, values)

        pptx_slide = getattr(slide, "pptx_slide", slide)

        graphic_frame = pptx_slide.shapes.add_chart(
            self._chart_type,
            Inches(self._left),
            Inches(self._top),
            Inches(self._width),
            Inches(self._height),
            chart_data,
        )

        chart = graphic_frame.chart

        if self._title:
            chart.has_title = True
            chart.chart_title.text_frame.text = self._title

        chart.has_legend = self._has_legend

        if self._has_data_labels:
            for plot in chart.plots:
                plot.has_data_labels = True

        return chart

    # ------------------------------------------------------------------
    # Convenience factory methods
    # ------------------------------------------------------------------

    @classmethod
    def from_dict(cls, spec: dict) -> "ChartBuilder":
        """Create a :class:`ChartBuilder` from a specification dictionary.

        The dictionary may contain the following keys (all optional unless
        noted):

        ``type`` *(required)*
            Chart type name (see :data:`CHART_TYPES`).
        ``categories``
            List of category label strings.
        ``series``
            List of ``{"name": str, "values": [num, ...]}`` dicts.
        ``title``
            Chart title string.
        ``legend``
            Boolean — show legend (default ``True``).
        ``data_labels``
            Boolean — show data labels (default ``False``).
        ``position``
            Dict with keys ``left``, ``top``, ``width``, ``height``
            (all in inches).
        """
        builder = cls(spec.get("type", "column"))

        if "categories" in spec:
            builder.categories(spec["categories"])

        for s in spec.get("series", []):
            builder.add_series(s["name"], s["values"])

        if "title" in spec:
            builder.title(spec["title"])

        if "legend" in spec:
            builder.legend(spec["legend"])

        if "data_labels" in spec:
            builder.data_labels(spec["data_labels"])

        if "position" in spec:
            p = spec["position"]
            builder.position(
                left=p.get("left", 1.0),
                top=p.get("top", 1.5),
                width=p.get("width", 8.0),
                height=p.get("height", 4.5),
            )

        return builder
