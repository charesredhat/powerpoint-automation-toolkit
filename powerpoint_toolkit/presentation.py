"""Presentation management — create, open, save, and organise slides."""

from __future__ import annotations

import copy
from pathlib import Path
from typing import List, Optional, Union

from pptx import Presentation as _PptxPresentation
from pptx.util import Inches

from powerpoint_toolkit.slides import Slide

# Slide-layout name → index mapping for the default theme.
# python-pptx exposes layouts by index; names vary per theme, but the
# standard "Office Theme" uses this ordering.
_LAYOUT_NAMES = {
    "blank": 6,
    "title": 0,
    "title_and_content": 1,
    "title_only": 5,
    "two_content": 3,
    "comparison": 4,
    "content_with_caption": 7,
    "picture_with_caption": 8,
    "section_header": 2,
}


class Presentation:
    """A high-level wrapper around a python-pptx *Presentation* object.

    Parameters
    ----------
    path:
        Path to an existing ``.pptx`` file to open.  When *None* a brand
        new, empty presentation is created.
    width_inches:
        Slide width in inches (default 13.33 — widescreen 16:9).
    height_inches:
        Slide height in inches (default 7.5 — widescreen 16:9).

    Examples
    --------
    Create a new presentation and add a title slide::

        prs = Presentation()
        slide = prs.add_slide(layout="title")
        slide.set_title("My Presentation")
        prs.save("output.pptx")

    Open an existing presentation::

        prs = Presentation("existing.pptx")
        print(prs.slide_count)
    """

    def __init__(
        self,
        path: Optional[Union[str, Path]] = None,
        *,
        width_inches: float = 13.33,
        height_inches: float = 7.5,
    ) -> None:
        if path is not None:
            self._prs = _PptxPresentation(str(path))
        else:
            self._prs = _PptxPresentation()
            self._prs.slide_width = Inches(width_inches)
            self._prs.slide_height = Inches(height_inches)

        self._path: Optional[Path] = Path(path) if path is not None else None

    # ------------------------------------------------------------------
    # Properties
    # ------------------------------------------------------------------

    @property
    def pptx(self) -> _PptxPresentation:
        """The underlying *python-pptx* :class:`~pptx.Presentation` object."""
        return self._prs

    @property
    def slide_count(self) -> int:
        """Number of slides in the presentation."""
        return len(self._prs.slides)

    @property
    def slides(self) -> List[Slide]:
        """All slides as :class:`~powerpoint_toolkit.slides.Slide` objects."""
        return [Slide(s, self._prs) for s in self._prs.slides]

    @property
    def slide_width(self) -> float:
        """Slide width in inches."""
        return self._prs.slide_width / 914400  # EMU → inches

    @property
    def slide_height(self) -> float:
        """Slide height in inches."""
        return self._prs.slide_height / 914400

    # ------------------------------------------------------------------
    # Slide management
    # ------------------------------------------------------------------

    def add_slide(self, layout: Union[str, int] = "blank") -> Slide:
        """Add a new slide and return a :class:`~powerpoint_toolkit.slides.Slide`.

        Parameters
        ----------
        layout:
            Either a layout name (``"blank"``, ``"title"``,
            ``"title_and_content"``, ``"title_only"``, ``"two_content"``,
            ``"comparison"``, ``"content_with_caption"``,
            ``"picture_with_caption"``, ``"section_header"``) or a
            zero-based integer index into the presentation's slide layout list.
        """
        layout_index = self._resolve_layout(layout)
        slide_layout = self._prs.slide_layouts[layout_index]
        pptx_slide = self._prs.slides.add_slide(slide_layout)
        return Slide(pptx_slide, self._prs)

    def get_slide(self, index: int) -> Slide:
        """Return the slide at *index* (zero-based).

        Raises
        ------
        IndexError
            When *index* is out of range.
        """
        slides = self._prs.slides
        if index < 0 or index >= len(slides):
            raise IndexError(
                f"Slide index {index} out of range (presentation has {len(slides)} slides)."
            )
        return Slide(slides[index], self._prs)

    def delete_slide(self, index: int) -> None:
        """Remove the slide at *index* (zero-based) from the presentation.

        Raises
        ------
        IndexError
            When *index* is out of range.
        """
        slides = self._prs.slides
        if index < 0 or index >= len(slides):
            raise IndexError(
                f"Slide index {index} out of range (presentation has {len(slides)} slides)."
            )
        xml_slides = self._prs.slides._sldIdLst  # noqa: SLF001
        xml_slides.remove(xml_slides[index])

    def duplicate_slide(self, index: int) -> Slide:
        """Duplicate the slide at *index* and append the copy to the end.

        Returns the newly created :class:`~powerpoint_toolkit.slides.Slide`.
        """
        source = self.get_slide(index)
        template_slide = source.pptx_slide

        # Add a blank slide of the same layout so python-pptx registers it
        blank_layout = self._prs.slide_layouts[_LAYOUT_NAMES["blank"]]
        new_pptx_slide = self._prs.slides.add_slide(blank_layout)

        # Replace the spTree (shape tree) with a deep copy of the source
        sp_tree = new_pptx_slide.shapes._spTree  # noqa: SLF001
        for child in list(sp_tree):
            sp_tree.remove(child)
        for child in template_slide.shapes._spTree:  # noqa: SLF001
            sp_tree.append(copy.deepcopy(child))

        return Slide(new_pptx_slide, self._prs)

    def reorder_slide(self, from_index: int, to_index: int) -> None:
        """Move the slide at *from_index* to *to_index*.

        Both indices are zero-based.  The operation is performed in-place.
        """
        slides_list = self._prs.slides._sldIdLst  # noqa: SLF001
        slide_count = len(slides_list)
        for idx in (from_index, to_index):
            if idx < 0 or idx >= slide_count:
                raise IndexError(
                    f"Slide index {idx} out of range (presentation has {slide_count} slides)."
                )
        element = slides_list[from_index]
        slides_list.remove(element)
        slides_list.insert(to_index, element)

    # ------------------------------------------------------------------
    # Persistence
    # ------------------------------------------------------------------

    def save(self, path: Optional[Union[str, Path]] = None) -> Path:
        """Save the presentation to *path*.

        When *path* is omitted the presentation is saved to the file it
        was originally loaded from.  If no path is available a
        :exc:`ValueError` is raised.

        Returns the resolved :class:`~pathlib.Path` that was written.
        """
        if path is None:
            if self._path is None:
                raise ValueError(
                    "No file path specified.  Call save(path) with an explicit path."
                )
            save_path = self._path
        else:
            save_path = Path(path)

        save_path.parent.mkdir(parents=True, exist_ok=True)
        self._prs.save(str(save_path))
        self._path = save_path
        return save_path

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _resolve_layout(self, layout: Union[str, int]) -> int:
        if isinstance(layout, int):
            return layout
        key = layout.lower().replace(" ", "_").replace("-", "_")
        if key not in _LAYOUT_NAMES:
            available = ", ".join(sorted(_LAYOUT_NAMES))
            raise ValueError(
                f"Unknown layout {layout!r}.  Available layouts: {available}"
            )
        index = _LAYOUT_NAMES[key]
        # Clamp to what this theme actually provides
        return min(index, len(self._prs.slide_layouts) - 1)

    def info(self) -> dict:
        """Return a dictionary with presentation metadata."""
        return {
            "slide_count": self.slide_count,
            "slide_width_inches": round(self.slide_width, 3),
            "slide_height_inches": round(self.slide_height, 3),
            "path": str(self._path) if self._path else None,
        }

    def __repr__(self) -> str:
        return (
            f"Presentation(path={self._path!r}, slides={self.slide_count})"
        )
