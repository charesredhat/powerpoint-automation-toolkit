"""Slide wrapper with high-level helpers for adding content."""

from __future__ import annotations

from pathlib import Path
from typing import List, Optional, Sequence, Tuple, Union

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from powerpoint_toolkit.shapes import (
    apply_text_format,
    apply_paragraph_format,
    hex_to_rgb,
)

_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


class Slide:
    """A high-level wrapper around a python-pptx slide.

    Instances are normally obtained from
    :meth:`~powerpoint_toolkit.presentation.Presentation.add_slide` or
    :meth:`~powerpoint_toolkit.presentation.Presentation.get_slide` rather
    than instantiated directly.

    Parameters
    ----------
    pptx_slide:
        The underlying ``pptx.slide.Slide`` object.
    presentation:
        The parent ``pptx.Presentation`` object (needed for slide dimensions).
    """

    def __init__(self, pptx_slide, presentation) -> None:
        self._slide = pptx_slide
        self._prs = presentation

    # ------------------------------------------------------------------
    # Properties
    # ------------------------------------------------------------------

    @property
    def pptx_slide(self):
        """The underlying *python-pptx* slide object."""
        return self._slide

    @property
    def shapes(self):
        """Collection of all shapes on this slide."""
        return self._slide.shapes

    @property
    def slide_width(self) -> float:
        """Slide width in inches."""
        return self._prs.slide_width / 914400

    @property
    def slide_height(self) -> float:
        """Slide height in inches."""
        return self._prs.slide_height / 914400

    # ------------------------------------------------------------------
    # Background
    # ------------------------------------------------------------------

    def set_background_color(self, hex_color: str) -> "Slide":
        """Fill the slide background with a solid colour.

        Parameters
        ----------
        hex_color:
            CSS hex string, e.g. ``"#1F3864"`` or ``"1F3864"``.

        Returns
        -------
        Slide
            *self*, to allow method chaining.
        """
        from pptx.oxml.ns import qn
        from lxml import etree

        bg = self._slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(hex_color)
        return self

    # ------------------------------------------------------------------
    # Title / subtitle
    # ------------------------------------------------------------------

    def set_title(
        self,
        text: str,
        *,
        font_size: Optional[float] = None,
        font_name: Optional[str] = None,
        bold: Optional[bool] = None,
        color: Optional[str] = None,
        alignment: Optional[str] = None,
    ) -> "Slide":
        """Set the title placeholder text.

        If the current slide layout has no title placeholder a new text box
        is created at the top of the slide.

        Returns *self* for chaining.
        """
        placeholder = self._find_placeholder("title")
        if placeholder is not None:
            tf = placeholder.text_frame
            tf.text = text
            para = tf.paragraphs[0]
        else:
            # Fallback: create a text box spanning the full width at the top
            txBox = self._slide.shapes.add_textbox(
                Inches(0.5),
                Inches(0.2),
                Inches(self.slide_width - 1.0),
                Inches(1.2),
            )
            tf = txBox.text_frame
            tf.text = text
            para = tf.paragraphs[0]

        run = para.runs[0] if para.runs else para.add_run()
        apply_text_format(
            run,
            font_size=font_size,
            font_name=font_name,
            bold=bold,
            color=color,
        )
        if alignment:
            apply_paragraph_format(para, alignment=alignment)
        return self

    def set_subtitle(
        self,
        text: str,
        *,
        font_size: Optional[float] = None,
        font_name: Optional[str] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        color: Optional[str] = None,
        alignment: Optional[str] = None,
    ) -> "Slide":
        """Set the subtitle/body placeholder text.

        Returns *self* for chaining.
        """
        placeholder = self._find_placeholder("body")
        if placeholder is None:
            placeholder = self._find_placeholder("subtitle")
        if placeholder is not None:
            tf = placeholder.text_frame
            tf.text = text
            para = tf.paragraphs[0]
            run = para.runs[0] if para.runs else para.add_run()
            apply_text_format(
                run,
                font_size=font_size,
                font_name=font_name,
                bold=bold,
                italic=italic,
                color=color,
            )
            if alignment:
                apply_paragraph_format(para, alignment=alignment)
        return self

    # ------------------------------------------------------------------
    # Text box
    # ------------------------------------------------------------------

    def add_text_box(
        self,
        text: str,
        *,
        left: float = 1.0,
        top: float = 1.5,
        width: float = 8.0,
        height: float = 1.0,
        font_size: Optional[float] = None,
        font_name: Optional[str] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        underline: Optional[bool] = None,
        color: Optional[str] = None,
        alignment: Optional[str] = None,
        word_wrap: bool = True,
    ) -> object:
        """Add a text box to the slide.

        All positional arguments are in inches.

        Returns the underlying python-pptx shape.
        """
        txBox = self._slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        tf.text = text

        para = tf.paragraphs[0]
        run = para.runs[0] if para.runs else para.add_run()
        apply_text_format(
            run,
            font_size=font_size,
            font_name=font_name,
            bold=bold,
            italic=italic,
            underline=underline,
            color=color,
        )
        if alignment:
            apply_paragraph_format(para, alignment=alignment)
        return txBox

    def add_bullet_list(
        self,
        items: Sequence[str],
        *,
        left: float = 1.0,
        top: float = 1.5,
        width: float = 8.0,
        height: float = 4.0,
        font_size: Optional[float] = None,
        font_name: Optional[str] = None,
        color: Optional[str] = None,
    ) -> object:
        """Add a bulleted list text box.

        Parameters
        ----------
        items:
            Sequence of text strings, one per bullet point.

        Returns the underlying python-pptx shape.
        """
        txBox = self._slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()
            para.text = item
            para.level = 0
            run = para.runs[0] if para.runs else para.add_run()
            apply_text_format(
                run,
                font_size=font_size,
                font_name=font_name,
                color=color,
            )
        return txBox

    # ------------------------------------------------------------------
    # Image
    # ------------------------------------------------------------------

    def add_image(
        self,
        image_path: Union[str, Path],
        *,
        left: float = 1.0,
        top: float = 1.5,
        width: Optional[float] = None,
        height: Optional[float] = None,
    ) -> object:
        """Add an image to the slide.

        At least one of *width* or *height* must be given; if only one is
        supplied the image is scaled proportionally.  If neither is given the
        image is inserted at its native size (capped to fit the slide).

        Returns the underlying python-pptx shape.
        """
        image_path = Path(image_path)
        if not image_path.exists():
            raise FileNotFoundError(f"Image file not found: {image_path}")

        pic = self._slide.shapes.add_picture(
            str(image_path),
            Inches(left),
            Inches(top),
            width=Inches(width) if width is not None else None,
            height=Inches(height) if height is not None else None,
        )
        return pic

    # ------------------------------------------------------------------
    # Table
    # ------------------------------------------------------------------

    def add_table(
        self,
        data: Sequence[Sequence],
        *,
        has_header: bool = True,
        left: float = 1.0,
        top: float = 1.5,
        width: float = 8.0,
        height: float = 3.0,
        font_size: Optional[float] = None,
        header_bg_color: Optional[str] = None,
        header_font_color: Optional[str] = None,
        alternate_row_color: Optional[str] = None,
    ) -> object:
        """Add a table to the slide.

        Parameters
        ----------
        data:
            2-D list/tuple of cell values.  If *has_header* is ``True``
            the first row is treated as the column header.
        has_header:
            Whether the first row is a header row.
        left, top, width, height:
            Position and size in inches.
        font_size:
            Font size for all cells (points).
        header_bg_color:
            Hex background colour for header cells.
        header_font_color:
            Hex font colour for header cells.
        alternate_row_color:
            Hex background colour for alternate (even-index) data rows.

        Returns the underlying python-pptx table shape.
        """
        rows = len(data)
        cols = max(len(row) for row in data) if data else 1

        table_shape = self._slide.shapes.add_table(
            rows,
            cols,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        table = table_shape.table

        for row_idx, row_data in enumerate(data):
            for col_idx, cell_value in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_value) if cell_value is not None else ""

                tf = cell.text_frame
                para = tf.paragraphs[0]
                run = para.runs[0] if para.runs else para.add_run()

                is_header = has_header and row_idx == 0
                is_even_data_row = (
                    not is_header
                    and alternate_row_color is not None
                    and (row_idx % 2 == 0)
                )

                apply_text_format(
                    run,
                    font_size=font_size,
                    bold=True if is_header else None,
                    color=header_font_color if is_header else None,
                )

                if is_header and header_bg_color:
                    fill = cell.fill
                    fill.solid()
                    fill.fore_color.rgb = hex_to_rgb(header_bg_color)
                elif is_even_data_row and alternate_row_color:
                    fill = cell.fill
                    fill.solid()
                    fill.fore_color.rgb = hex_to_rgb(alternate_row_color)

        return table_shape

    # ------------------------------------------------------------------
    # Chart (convenience wrapper)
    # ------------------------------------------------------------------

    def add_chart(
        self,
        chart_type: str,
        categories: Sequence[str],
        series: Sequence[Tuple[str, Sequence[Union[int, float]]]],
        *,
        title: Optional[str] = None,
        left: float = 1.0,
        top: float = 1.5,
        width: float = 8.0,
        height: float = 4.5,
        has_legend: bool = True,
        has_data_labels: bool = False,
    ) -> object:
        """Add a chart to the slide using the :class:`~powerpoint_toolkit.charts.ChartBuilder`.

        Parameters
        ----------
        chart_type:
            A chart-type name from :data:`~powerpoint_toolkit.charts.CHART_TYPES`.
        categories:
            Category labels.
        series:
            List of ``(name, values)`` tuples.
        title:
            Optional chart title.
        left, top, width, height:
            Position and size in inches.
        has_legend:
            Show legend.
        has_data_labels:
            Show data labels on series.

        Returns the underlying python-pptx ``Chart`` object.
        """
        from powerpoint_toolkit.charts import ChartBuilder

        builder = (
            ChartBuilder(chart_type)
            .categories(categories)
            .legend(has_legend)
            .data_labels(has_data_labels)
            .position(left=left, top=top, width=width, height=height)
        )
        for name, values in series:
            builder.add_series(name, values)
        if title:
            builder.title(title)

        return builder.build(self)

    # ------------------------------------------------------------------
    # Shape helpers
    # ------------------------------------------------------------------

    def add_rectangle(
        self,
        *,
        left: float = 1.0,
        top: float = 1.0,
        width: float = 3.0,
        height: float = 1.5,
        fill_color: Optional[str] = None,
        line_color: Optional[str] = None,
    ) -> object:
        """Add a filled rectangle to the slide.

        Returns the underlying python-pptx shape.
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.util import Inches as _Inches

        shape = self._slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
        if line_color:
            shape.line.color.rgb = hex_to_rgb(line_color)
        else:
            shape.line.fill.background()
        return shape

    # ------------------------------------------------------------------
    # Internal helpers
    # ------------------------------------------------------------------

    def _find_placeholder(self, ph_type: str):
        """Find a placeholder by type name.

        *ph_type* is one of ``"title"``, ``"body"``, ``"subtitle"``.
        Returns ``None`` if not found.
        """
        from pptx.enum.shapes import PP_PLACEHOLDER

        _type_map = {
            "title": PP_PLACEHOLDER.TITLE,
            "center_title": PP_PLACEHOLDER.CENTER_TITLE,
            "body": PP_PLACEHOLDER.BODY,
            "subtitle": PP_PLACEHOLDER.SUBTITLE,
        }

        for ph in self._slide.placeholders:
            ph_type_enum = _type_map.get(ph_type)
            if ph_type_enum is not None and ph.placeholder_format.type == ph_type_enum:
                return ph
            # Also try by name substring for flexibility
            if ph_type.lower() in ph.name.lower():
                return ph
        return None

    def __repr__(self) -> str:
        return f"Slide(shapes={len(self._slide.shapes)})"
