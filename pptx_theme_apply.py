from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
import shutil
from io import BytesIO
import os

SRC = r"C:\Project\Powerpoint\HFP\CDER Research Governance WOAP presentation.pptx"
TGT = r"C:\Project\Powerpoint\HFP\HPC_Cloud_Analysis_HPC_GAB_Consolidated_Slides_2.pptx"
BACKUP = TGT.replace('.pptx', '_backup.pptx')
OUT = TGT.replace('.pptx', '_theme_applied.pptx')
EXTRACT = os.path.join(os.path.dirname(TGT), 'HPC_extracted_text.txt')

shutil.copy2(TGT, BACKUP)
print('Backup created:', BACKUP)

src = Presentation(SRC)
tgt = Presentation(TGT)

# Try to get a representative background image from source (first picture found)
bg_blob = None
for slide in src.slides:
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                bg_blob = shape.image.blob
                break
            except Exception:
                continue
    if bg_blob:
        break

# Try to get a representative font name + size from source (first run found)
font_name = None
font_size = None
for slide in src.slides:
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame is not None:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.name:
                        font_name = run.font.name
                    if run.font.size:
                        font_size = run.font.size
                    if font_name or font_size:
                        break
                if font_name or font_size:
                    break
        if font_name or font_size:
            break
    if font_name or font_size:
        break

print('Detected font:', font_name, 'size:', font_size)

# Extract all text to a file
with open(EXTRACT, 'w', encoding='utf-8') as f:
    for i, slide in enumerate(tgt.slides, start=1):
        f.write(f"--- Slide {i} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text = shape.text.strip()
                if text:
                    f.write(text + "\n")
        f.write('\n')

print('Extracted text written to', EXTRACT)

# Apply background and fonts (best-effort)
from pptx.util import Pt

for slide in tgt.slides:
    # Apply background image if found
    if bg_blob:
        try:
            img_stream = BytesIO(bg_blob)
            slide.shapes.add_picture(img_stream, 0, 0, width=tgt.slide_width, height=tgt.slide_height)
        except Exception as e:
            print('Failed to add background image to a slide:', e)
    # Apply font changes
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame') or shape.text_frame is None:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if font_name:
                    try:
                        run.font.name = font_name
                    except Exception:
                        pass
                if font_size:
                    try:
                        run.font.size = font_size
                    except Exception:
                        pass

# Save output
try:
    tgt.save(OUT)
    print('Saved themed presentation as', OUT)
except Exception as e:
    print('Failed to save edited presentation:', e)

print('DONE')
